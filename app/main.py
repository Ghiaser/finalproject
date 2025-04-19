import os
import base64
import hashlib
import numpy as np
from typing import List
from PIL import Image, ImageSequence

import torch
from transformers import CLIPProcessor, CLIPModel
import faiss
from cryptography.fernet import Fernet
import pickle

from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation

def generate_fernet_key(password: str) -> bytes:
    hashed = hashlib.sha256(password.encode()).digest()
    return base64.urlsafe_b64encode(hashed)

def encrypt_vector(vec: np.ndarray, fernet: Fernet) -> bytes:
    return fernet.encrypt(vec.tobytes())

def decrypt_vector(enc_bytes: bytes, shape, fernet: Fernet) -> np.ndarray:
    return np.frombuffer(fernet.decrypt(enc_bytes), dtype=np.float32).reshape(shape)

def extract_text_from_pdf(path):
    try:
        reader = PdfReader(path)
        return "\n".join([page.extract_text() or "" for page in reader.pages])
    except:
        return ""

def extract_text_from_docx(path):
    try:
        doc = Document(path)
        return "\n".join([p.text for p in doc.paragraphs])
    except:
        return ""

def extract_text_from_pptx(path):
    try:
        prs = Presentation(path)
        return "\n".join([shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text")])
    except:
        return ""

def extract_image_from_gif(path):
    try:
        img = Image.open(path)
        return next(ImageSequence.Iterator(img)).convert("RGB")
    except:
        return None

class CLIPSecureEncryptor:
    def __init__(self, password: str, device='cuda' if torch.cuda.is_available() else 'cpu'):
        self.device = device
        self.processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
        self.model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32").to(self.device)
        self.fernet = Fernet(generate_fernet_key(password))
        self.index = None
        self.data_refs = []
        self.encrypted_vectors = []

    def encode_file(self, path: str) -> np.ndarray:
        ext = os.path.splitext(path)[1].lower()
        content = ""

        if ext in [".jpg", ".jpeg", ".png"]:
            image = Image.open(path).convert("RGB")
            inputs = self.processor(images=image, return_tensors="pt").to(self.device)
            with torch.no_grad():
                return self.model.get_image_features(**inputs).cpu().numpy().astype("float32")

        elif ext == ".gif":
            image = extract_image_from_gif(path)
            if image:
                inputs = self.processor(images=image, return_tensors="pt").to(self.device)
                with torch.no_grad():
                    return self.model.get_image_features(**inputs).cpu().numpy().astype("float32")

        elif ext == ".pdf":
            content = extract_text_from_pdf(path)
        elif ext == ".docx":
            content = extract_text_from_docx(path)
        elif ext == ".pptx":
            content = extract_text_from_pptx(path)
        else:
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = f.read()
            except:
                pass

        inputs = self.processor(text=content, return_tensors="pt", padding=True, truncation=True, max_length=77).to(self.device)
        with torch.no_grad():
            return self.model.get_text_features(**inputs).cpu().numpy().astype("float32")

    def build_index_from_files(self, file_paths: List[str]):
        self.model.eval()
        self.encrypted_vectors.clear()
        self.data_refs.clear()

        for path in file_paths:
            try:
                vec = self.encode_file(path)
                enc_vec = encrypt_vector(vec, self.fernet)
                self.encrypted_vectors.append(enc_vec)
                self.data_refs.append(self.fernet.encrypt(path.encode()).decode())
            except Exception as e:
                print(f"Error encoding {path}: {e}")

        if self.encrypted_vectors:
            decrypted_vectors = [decrypt_vector(enc, (-1, vec.shape[1]), self.fernet) for enc in self.encrypted_vectors]
            matrix = np.vstack(decrypted_vectors).astype("float32")
            self.index = faiss.IndexFlatL2(matrix.shape[1])
            self.index.add(matrix)

    def _decrypt_refs(self):
        return [self.fernet.decrypt(ref.encode()).decode() for ref in self.data_refs]

    def query_text(self, text: str, k=3):
        inputs = self.processor(text=text, return_tensors="pt", padding=True, truncation=True, max_length=77).to(self.device)
        with torch.no_grad():
            emb = self.model.get_text_features(**inputs).cpu().numpy().astype("float32")
        enc_vec = encrypt_vector(emb, self.fernet)
        query_vec = decrypt_vector(enc_vec, emb.shape, self.fernet)
        distances, indices = self.index.search(query_vec, k)
        return [self._decrypt_refs()[i] for i in indices[0]]

    def query_file(self, path: str, k=3):
        vec = self.encode_file(path)
        enc_vec = encrypt_vector(vec, self.fernet)
        query_vec = decrypt_vector(enc_vec, vec.shape, self.fernet)
        distances, indices = self.index.search(query_vec, k)
        return [self._decrypt_refs()[i] for i in indices[0]]

    def save_index(self, path="secure_index.pkl"):
        with open(path, "wb") as f:
            pickle.dump({"encrypted_vectors": self.encrypted_vectors, "data_refs": self.data_refs}, f)

    def load_index(self, path="secure_index.pkl"):
        with open(path, "rb") as f:
            obj = pickle.load(f)
            self.encrypted_vectors = obj["encrypted_vectors"]
            self.data_refs = obj["data_refs"]

        if self.encrypted_vectors:
            decrypted_vectors = [decrypt_vector(enc, (-1, 512), self.fernet) for enc in self.encrypted_vectors]  # Assuming 512 dim
            matrix = np.vstack(decrypted_vectors).astype("float32")
            self.index = faiss.IndexFlatL2(matrix.shape[1])
            self.index.add(matrix)
